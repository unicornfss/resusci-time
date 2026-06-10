"""Generate Resusci-Time QI briefing PowerPoint (compact, on-brand)."""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "Resusci-Time-QI-and-Technical-Overview.pptx"
OUTPUT_FALLBACK = ROOT / "Resusci-Time-QI-and-Technical-Overview-new.pptx"
WMAS_CREST = ROOT / "public" / "backgrounds" / "wmas-crest.png"
RESUSCI_LOGO = ROOT / "public" / "backgrounds" / "resusci-time-logo.png"

# Resusci-Time day theme (index.css / manifest)
GREEN_800 = RGBColor(0x16, 0x38, 0x16)
GREEN_700 = RGBColor(0x1F, 0x4F, 0x1F)
GREEN_600 = RGBColor(0x2D, 0x6A, 0x2D)
GREEN_100 = RGBColor(0xD0, 0xE0, 0xD0)
BG = RGBColor(0xDF, 0xE8, 0xDF)
TEXT = RGBColor(0x1A, 0x2E, 0x1A)
MUTED = RGBColor(0x4A, 0x5F, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def load_app_version() -> str:
    package_json = json.loads((ROOT / "package.json").read_text(encoding="utf-8"))
    return package_json["version"]


def load_preview_date() -> str:
    changelog = (ROOT / "TESTING-CHANGELOG.md").read_text(encoding="utf-8")
    match = re.search(r"\*\*Last updated:\*\*\s*([^\n]+)", changelog)
    if match:
        return match.group(1).strip()
    return "See TESTING-CHANGELOG.md"


APP_VERSION = load_app_version()
PREVIEW_DATE = load_preview_date()

# Version | Released | Summary (condensed from TESTING-CHANGELOG.md and release history)
VERSION_CHANGELOG = [
    [
        "1.0.0",
        "27 May 2026",
        "Initial release — adult arrest timer and checklist; VF/pVT, PEA, Asystole; rhythm checks, adrenaline, amiodarone; ROSC, TOR, VoD; interventions; 4 Hs / 4 Ts; offline installable PWA",
    ],
    [
        "1.0.1",
        "28 May 2026",
        "Multi-trust builds (Standard / WMAS); event log CSV/PDF export, share link and QR; saved logs; public blog; trust branding and separate deploy URLs",
    ],
    [
        "1.1.0",
        "28 May 2026",
        "Screen wake lock during active case; confirm before starting new case; blog audience filter by trust",
    ],
    [
        "1.1.1",
        "28–31 May 2026",
        "Documents modal (ALS algorithm); prolonged VF reminder; TOR special circumstances and reassessment; VoD observation flow; CODE SHOCK and WMAS ToR document (WMAS); preview speed controls; needle decompression",
    ],
    [
        "1.1.2",
        "2 Jun 2026",
        "Clinical alerts one at a time; autosave logs and multi-delete; continue case within 10 minutes; DEMO preview icons; preview startup warning",
    ],
    [
        "1.2.0",
        "2 Jun 2026",
        "Case transfer — pause timer, QR handoff to another device, sender read-only when complete",
    ],
    [
        "1.2.1",
        "2 Jun 2026",
        "Transfer warning if check or adrenaline due soon; Cardiac arrest button from ROSC (re-arrest); sustained ROSC alert; SBP/pulse reminder fixes",
    ],
    [
        "1.2.2",
        "4 Jun 2026",
        "WMAS CODE SHOCK only when initial rhythm was VF / pVT (not after non-shockable initial rhythm)",
    ],
    [
        "1.2.3",
        "6 Jun 2026",
        "Patient handed over flow; timer bar layout and colours; hamburger menu; acknowledgements page; About updates",
    ],
    [
        "1.2.4",
        "7 Jun 2026",
        "Preview debug report export (menu); session timeline and error capture for testers",
    ],
    [
        "1.2.5",
        "9 Jun 2026",
        "Metronome stops on ROSC but resumes on re-arrest if still on; preview changelog and landing-page version deploy fixes",
    ],
    [
        "1.2.6",
        "9 Jun 2026",
        "Case transfer (QR) custom/WMAS only — removed from Standard build; Patient handed over remains on all builds",
    ],
    [
        "1.2.7",
        "10 Jun 2026",
        "Phone timer-bar layout; scroll-for-checklist hint; case continuation uses protocol minutes and last log entry; Standard handover copy",
    ],
]

# Preserved from user-edited deck — do not change without checking the .pptx first.
REPLACE_SOMETHING_TABLE = {
    "title": "3. Are you trying to replace something?",
    "headers": ["Today", "Resusci-Time"],
    "rows": [
        ["JRCALC / paper protocols", "Same guidance — adds active timer + log"],
        ["Informal timing / mental tracking", "Supplements, does not remove team leader role"],
        ["EPR / PRF", "Export may help after the event — not a replacement"],
        ["Another mandated arrest app", "Only relevant if trust already mandates one"],
    ],
    "note": "Optional aide for timing, prompts, and logging — not a swap for authoritative sources.",
}


def set_run_font(run, size=16, bold=False, color=TEXT):
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_brand_background(slide, prs):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_700
    bar.line.fill.background()


def add_title_slide(prs, title, subtitle="", show_logos=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    panel = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = GREEN_700
    panel.line.fill.background()

    if show_logos:
        if WMAS_CREST.exists():
            slide.shapes.add_picture(str(WMAS_CREST), Inches(0.85), Inches(0.55), height=Inches(1.25))
        if RESUSCI_LOGO.exists():
            slide.shapes.add_picture(str(RESUSCI_LOGO), Inches(10.55), Inches(0.65), height=Inches(0.95))

    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.05), Inches(11.5), Inches(1.2))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run_font(run, size=36, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.85), Inches(3.2), Inches(11.2), Inches(2.4))
        tf = sub.text_frame
        tf.word_wrap = True
        for i, line in enumerate(subtitle.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_after = Pt(6)
            sr = para.add_run()
            sr.text = line
            set_run_font(sr, size=17, color=GREEN_100)


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(slide, prs)
    box = slide.shapes.add_textbox(Inches(0.85), Inches(2.9), Inches(11.5), Inches(1.2))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = title
    set_run_font(run, size=32, bold=True, color=GREEN_800)
    underline = slide.shapes.add_shape(1, Inches(0.85), Inches(3.55), Inches(2.2), Inches(0.07))
    underline.fill.solid()
    underline.fill.fore_color.rgb = GREEN_600
    underline.line.fill.background()


def add_bullet_slide(prs, title, bullets, note=None, title_size=26, body_size=15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(slide, prs)

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.7))
    tr = title_box.text_frame.paragraphs[0].add_run()
    tr.text = title
    set_run_font(tr, size=title_size, bold=True, color=GREEN_800)

    body = slide.shapes.add_textbox(Inches(0.95), Inches(1.35), Inches(11.3), Inches(5.4))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = text
        set_run_font(run, size=body_size if level == 0 else 13, color=TEXT if level == 0 else MUTED)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.95), Inches(6.45), Inches(11.3), Inches(0.55))
        nr = nb.text_frame.paragraphs[0].add_run()
        nr.text = note
        set_run_font(nr, size=12, color=MUTED)


def add_compact_table_slide(prs, title, headers, rows, note=None, header_size=12, body_size=11):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_background(slide, prs)

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.7))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = title
    set_run_font(title_run, size=26, bold=True, color=GREEN_800)

    cols = len(headers)
    row_count = len(rows) + 1
    height = min(Inches(0.4 * row_count), Inches(5.2))
    table = slide.shapes.add_table(row_count, cols, Inches(0.85), Inches(1.3), Inches(11.4), height).table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN_600
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_run_font(run, size=header_size, bold=True, color=WHITE)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size=body_size, color=TEXT)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.85), Inches(6.45), Inches(11.4), Inches(0.55))
        nr = nb.text_frame.paragraphs[0].add_run()
        nr.text = note
        set_run_font(nr, size=12, color=MUTED)


def add_changelog_table_slide(prs, title, rows, note=None):
    add_compact_table_slide(
        prs,
        title,
        ["Version", "Released", "Summary"],
        rows,
        note=note,
        header_size=11,
        body_size=9,
    )


def add_version_history_slides(prs):
    add_section_slide(prs, "Version history")

    chunk_size = 4
    chunks = [
        VERSION_CHANGELOG[i : i + chunk_size]
        for i in range(0, len(VERSION_CHANGELOG), chunk_size)
    ]
    total = len(chunks)
    for index, chunk in enumerate(chunks, start=1):
        title = f"Changelog from v1.0.0 ({index} of {total})"
        note = None
        if index == total:
            note = f"Current preview build: v{APP_VERSION} · {PREVIEW_DATE}. Live field builds on main may trail preview until merged."
        add_changelog_table_slide(prs, title, chunk, note=note)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Resusci-Time — WMAS",
        f"Quality improvement briefing · Preview build v{APP_VERSION}\n{PREVIEW_DATE}\nJon Ostrowski · Paramedic / Clinical Team Mentor",
        show_logos=True,
    )

    add_bullet_slide(
        prs,
        "App features — case flow & timer",
        [
            "Start protocol, initial assessment, and initial rhythm (VF / pVT, PEA, Asystole)",
            "Elapsed resuscitation timer with 2-minute rhythm-check countdown",
            "Shock logging (joules), total shock count, adrenaline / amiodarone due in timer bar",
            "Switch to ROSC mode or return to cardiac arrest (re-arrest) from timer bar",
            "TOR button, 45-minute termination review alert, and full TOR questionnaire",
            "Post-TOR verification-of-death wait and initial-assessment VoD (asystole observation)",
            "Patient handed over — stops timers, locks case, opens log (all builds)",
            "Metronome on timer bar; night mode; screen stays awake during a case",
        ],
        body_size=13,
    )

    add_bullet_slide(
        prs,
        "App features — checklists, prompts & decisions",
        [
            "Resuscitation quality checklist and 4 Hs / 4 Ts reversible causes",
            "ROSC checklist, SBP / pulse monitoring prompts, atropine tracking (ROSC)",
            "Sustained ROSC alert (more than 10 minutes) and senior-discussion notice",
            "Interventions: vascular access, airway, breathing, circulation, medications",
            "Vector-change reminder after three consecutive shockable rhythms",
            "Early transfer reminder on third rhythm check (VF / pVT or PEA initial rhythm)",
            ("WMAS: CODE SHOCK to EOC after first shock (VF / pVT initial rhythm only)", 1),
            ("WMAS: Prolonged VF reminder; prolonged VF gates TOR senior discussion", 1),
            "Clinical alerts shown one at a time in priority order",
            "Vascular-access prompt if first adrenaline given without IV/IO logged",
        ],
        body_size=13,
    )

    add_bullet_slide(
        prs,
        "App features — log, transfer & tools",
        [
            "Timestamped event log; view during case; export CSV or PDF",
            "Autosave on device; saved logs list; continue case within 10 minutes of last entry",
            ("WMAS / custom: Case transfer — pause timer, QR to another device, sender read-only when complete", 1),
            ("WMAS / custom: Transfer warning if rhythm check or adrenaline due within one minute", 1),
            "Documents: ALS algorithm; WMAS ToR criteria aide-memoire (WMAS build)",
            "Menu: About, Documents, Saved logs, Acknowledgements, Install app",
            "During a case: header shows Night mode; WMAS / custom builds also show Transfer case",
            ("Preview only: debug report export, speed control (1×–10×), DEMO icons, startup warning", 1),
        ],
        note=f"Preview build v{APP_VERSION} — features on preview URLs until merged to main.",
        body_size=13,
    )

    add_version_history_slides(prs)

    add_section_slide(prs, "Questions from the QI lead")

    add_bullet_slide(
        prs,
        "1. What is its purpose?",
        [
            "Guided timer + checklist for adult cardiac arrest (pre-hospital)",
            "Prompts rhythm checks, drugs, ROSC monitoring, TOR, and VoD in line with JRCALC / AACE",
            "Timestamped event log for handover, debrief, and QA",
            "WMAS build adds trust-specific prompts (e.g. CODE SHOCK to EOC)",
            "Support tool only — not a substitute for JRCALC, trust policy, or clinical judgement",
        ],
        body_size=14,
    )

    add_compact_table_slide(
        prs,
        "2. What are you hoping to achieve?",
        ["Goal", "Why"],
        [
            ["Fewer missed rhythm checks / drug intervals", "Offloads time-keeping from working memory"],
            ["Consistent contemporaneous log", "Handover, debrief, governance"],
            ["Structured end-of-case flows", "ROSC, TOR, VoD, handover / transfer"],
            ["Low-friction pilot evaluation", "PWA on phone or tablet — no app store"],
        ],
    )

    add_compact_table_slide(
        prs,
        REPLACE_SOMETHING_TABLE["title"],
        REPLACE_SOMETHING_TABLE["headers"],
        REPLACE_SOMETHING_TABLE["rows"],
        note=REPLACE_SOMETHING_TABLE["note"],
    )

    add_compact_table_slide(
        prs,
        "4. Cognitive load — reduce or increase?",
        ["Reduces", "Risks if misused"],
        [
            ["One shared clock + due prompts", "Extra screen during hands-on care"],
            ["Externalised checklists (4 Hs / 4 Ts, ROSC)", "Logging takes seconds — scribe role"],
            ["Less mental arithmetic on drug timing", "Wrong initial rhythm → wrong branch"],
            ["Shared reference if crew changes", "Principle: one person runs the app"],
        ],
    )

    add_compact_table_slide(
        prs,
        "5. Unintended consequences?",
        ["Risk", "Mitigation"],
        [
            ["Over-reliance on prompts", "Training: supports decisions, does not make them"],
            ["Distraction from patient", "Designate scribe; install before shift"],
            ["Informal adoption without oversight", "Controlled pilot with QI / clinical lead"],
            ["Log treated as sole record", "Export supplements EPR — device-only by default"],
            ["Out-of-date build / protocol drift", "Versioning, preview testing, clinical review"],
        ],
    )

    add_bullet_slide(
        prs,
        "Suggested 30-second opener",
        [
            "Resusci-Time is a browser-based timer and checklist for adult cardiac arrest, aligned to Spring 2026 JRCALC / AACE guidance.",
            "It reduces the burden of tracking cycles, drugs, and documentation during a stressful job.",
            "It does not replace JRCALC, trust policy, or clinical judgement.",
            "The WMAS preview build includes CODE SHOCK and transfer / handover flows — I am seeking views on a controlled pilot and guardrails.",
        ],
        body_size=14,
    )

    add_compact_table_slide(
        prs,
        "How it is built",
        ["Part", "Role in the app", "Real-life analogy"],
        [
            [
                "TypeScript + React",
                "Draws screens and responds when you tap buttons",
                "Like the controls and display on a defib — what the crew sees and presses",
            ],
            [
                "Vite",
                "Packages the app into files a browser can load",
                "Like assembling the protocol pack into something crews can open on shift",
            ],
            [
                "PWA (installable web app)",
                "Add to home screen — opens full-screen like a normal app",
                "Like pinning the ALS algorithm to your tablet home screen",
            ],
            [
                "Service worker",
                "Keeps a copy on the device for offline use after first visit",
                "Like keeping a laminated protocol in the ambulance — still there with no signal",
            ],
            [
                "IndexedDB",
                "Stores saved logs on this device only",
                "Like a notebook that stays in the vehicle — not sent to a server",
            ],
            [
                "No backend server",
                "During a case nothing is uploaded",
                "Like a paper PRF that never leaves the crew’s hands until they choose to export",
            ],
            [
                "jsPDF",
                "Turns the event log into a PDF",
                "Like printing the handwritten event log for handover or QA",
            ],
            [
                "QR code (transfer)",
                "Encodes the live case for another device to scan (custom builds)",
                "Like photocopying the running log for the second crew — but live and digital",
            ],
            [
                "GitHub Pages + Cloudflare",
                "Hosts the app at a fixed web address",
                "Like the trust link on the intranet where crews install the right version",
            ],
            [
                "testing → preview / main → live",
                "Preview for testers; main for approved field builds",
                "Like draft SOP on SharePoint → signed-off version on the official site",
            ],
        ],
        note=f"Preview version {APP_VERSION} · jon@ostroforge.co.uk",
        header_size=10,
        body_size=9,
    )

    add_compact_table_slide(
        prs,
        "How individual trust versions are made",
        ["Layer", "Role", "Real-life analogy"],
        [
            [
                "One shared app",
                "Same core timer, checklist, log, and flows for all trusts",
                "Like one national JRCALC arrest chapter — one source of truth",
            ],
            [
                "Trust config file",
                "Per-trust name, crest, prompts (e.g. CODE SHOCK), extra documents",
                "Like the WMAS crest and local drug chart stapled into the front of the pack",
            ],
            [
                "Custom-only features",
                "Case transfer, CODE SHOCK, timing overrides — Standard excludes QR transfer",
                "Like optional trust inserts (e.g. EOC call-out) not every ambulance service uses",
            ],
            [
                "If not specified → defaults",
                "Rhythm check 2 min, adrenaline 4 min unless overridden",
                "Like JRCALC default intervals unless local SOP says otherwise",
            ],
            [
                "Optional overrides",
                "Trust-specific intervals or rules read at build time",
                "Like agreeing a 90-second rhythm check in config instead of changing the whole app",
            ],
            [
                "Separate build & URL",
                "WMAS vs Standard compiled separately — crews get the right install link",
                "Like different QR codes on WMAS vs neighbouring trust tablets",
            ],
            [
                "Preview → live",
                "Test on preview URL first; approved changes go to the field build",
                "Like pilot on training iPads before rolling out to front-line devices",
            ],
        ],
        note="Clinical changes need trust sign-off; most trust differences are configuration, not a rewrite of the app.",
        header_size=10,
        body_size=9,
    )

    add_compact_table_slide(
        prs,
        "Fixing simple errors and omissions",
        ["Type of fix", "Example", "Typical turnaround"],
        [
            [
                "Trust-only addition",
                "A drug missing from Interventions → Medications for one trust",
                "Same day on preview — add to trust config, rebuild, deploy (~30–60 min)",
            ],
            [
                "Shared list update",
                "A drug that should appear in all builds",
                "Same day on preview — one code change + quick test (~1–2 hours)",
            ],
            [
                "Wording or label",
                "Typo in a prompt, button, or log entry text",
                "Same day on preview (~15–30 min)",
            ],
            [
                "Trust reminder rule",
                "Adjust when a prompt appears (e.g. CODE SHOCK timing)",
                "Same day to 1 day on preview (~1–2 hours)",
            ],
            [
                "Live field build",
                "Any of the above after clinical review",
                "After merge to main and redeploy — usually within days, not weeks",
            ],
            [
                "Larger change",
                "New flow, timer logic, or clinical branch",
                "Longer — scoped separately; not a same-day fix",
            ],
        ],
        note="Preview URL can be updated quickly for testers; live crew builds follow trust sign-off. Crews may need to refresh or reinstall to pick up updates.",
        body_size=10,
    )

    add_section_slide(prs, "Future direction")

    add_compact_table_slide(
        prs,
        "Trust-only iOS / Android (not public app stores)",
        ["How it works", "What it means for Resusci-Time"],
        [
            [
                "Personally issued iPads enrolled in trust MDM (e.g. Microsoft Intune)",
                "Crew install from Company Portal — not by searching the App Store",
            ],
            [
                "IT assigns apps to WMAS device groups only",
                "App is private to the trust; updates controlled centrally",
            ],
            [
                "iOS: Apple Business Manager + custom / unlisted app",
                "Same React codebase can be wrapped (e.g. Capacitor) for distribution",
            ],
            [
                "Android: Managed Google Play private channel",
                "Crew iPads and CFR phones can each get the right build",
            ],
        ],
        note="Today: installable web app (PWA). Next step: wrapped app in Company Portal when IT and governance agree.",
        body_size=10,
    )

    add_bullet_slide(
        prs,
        "CFR version — limited scope, upgrade on handover",
        [
            "Separate CFR build: CPR quality, AED / shocks, rhythm timing, call for help — limited to CFR scope of practice",
            "Prompts match CFR scope of practice; buttons for IV drugs, advanced airway, TOR, etc. are hidden",
            "When ambulance arrives: Transfer case (QR) — same log and timers continue on the crew device",
            "Crew app unlocks full ALS prompts; CFR actions stay in the log; event logged as scope upgraded to crew",
            "Reuses existing case-transfer design — extend config with pathway (CFR vs crew) in the handoff payload",
            "Distribution: separate Company Portal entry for CFR devices, subject to trust IT and CFR governance",
        ],
        body_size=13,
    )

    add_bullet_slide(
        prs,
        "Paediatric version — separate clinical product",
        [
            "Not a toggle on the adult app — separate build (e.g. Resusci-Time Paediatric — WMAS) to avoid wrong-pathway use",
            "Age / weight bands, paediatric drug doses, shock energy, and algorithm branches from JRCALC paediatric guidance",
            "Same platform pattern as WMAS vs Standard: trust branding, documents, and reminder rules in config",
            "Dedicated clinical review (trust paediatric lead / PICU liaison) before any field use",
            "Likely follows adult crew pilot and CFR pathway unless trust has a specific paediatric QI priority",
        ],
        body_size=13,
    )

    add_bullet_slide(
        prs,
        "Event log filters (planned)",
        [
            "Filter the event log during review, debrief, or handover — full log stays intact; filters only change what is shown",
            "Suggested filters: Rhythm — rhythm checks, ROSC, cardiac arrest, and related entries",
            "Drugs — adrenaline, amiodarone, atropine, and other medication lines",
            "Shocks — defibrillation energies and shockable rhythm entries",
            "Interventions — airway, breathing, circulation, vascular access, and other logged actions",
            "Useful for quick handover summaries and post-case review without scrolling the entire timeline",
        ],
        note="Not in the current preview build — a future enhancement to the log viewer and export.",
        body_size=13,
    )

    add_title_slide(prs, "Thank you", "Questions and feedback welcome\njon@ostroforge.co.uk")

    target = OUTPUT
    try:
        prs.save(target)
    except PermissionError:
        target = OUTPUT_FALLBACK
        prs.save(target)
        print(f"Note: close the open .pptx and re-run to overwrite {OUTPUT.name}")
    print(f"Created: {target} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
